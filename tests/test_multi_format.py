"""
Tests for multi-format text extraction — .pptx and .xlsx support.

Creates disposable test files on the fly (no real files needed).
Also tests the dispatcher ``extract_text()`` and ``detect_file_type()``.
"""
import os

import openpyxl
import pytest
from pptx import Presentation
from pptx.util import Inches

from app.services.document_service import (
    extract_text,
    extract_text_from_pptx,
    extract_text_from_xlsx,
)
from app.utils.file_utils import detect_file_type


# ---------------------------------------------------------------------------
# Helpers — create disposable test files
# ---------------------------------------------------------------------------

def _create_pptx(tmp_dir: str, slides: list[list[str]], filename: str = "test.pptx") -> str:
    """Create a .pptx with given slide contents.

    ``slides`` is a list of lists: each inner list contains text paragraphs for one slide.
    """
    path = os.path.join(tmp_dir, filename)
    prs = Presentation()
    for slide_texts in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        for i, text in enumerate(slide_texts):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i * 0.5), Inches(4), Inches(0.4))
            txBox.text_frame.text = text
    prs.save(path)
    return path


def _create_pptx_with_table(tmp_dir: str, table_data: list[list[str]]) -> str:
    """Create a .pptx with a single slide containing a table."""
    path = os.path.join(tmp_dir, "table.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rows, cols = len(table_data), len(table_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(5), Inches(2))
    table = table_shape.table
    for r, row_data in enumerate(table_data):
        for c, cell_text in enumerate(row_data):
            table.cell(r, c).text = cell_text

    prs.save(path)
    return path


def _create_xlsx(
    tmp_dir: str,
    sheets: dict[str, list[list]],
    filename: str = "test.xlsx",
) -> str:
    """Create a .xlsx with named sheets and cell data.

    ``sheets`` maps sheet name → list of rows (each row is a list of cell values).
    """
    path = os.path.join(tmp_dir, filename)
    wb = openpyxl.Workbook()
    # Remove default sheet
    wb.remove(wb.active)

    for sheet_name, rows in sheets.items():
        ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)

    wb.save(path)
    return path


def _create_xlsx_empty(tmp_dir: str) -> str:
    """Create a .xlsx with an empty sheet."""
    path = os.path.join(tmp_dir, "empty.xlsx")
    wb = openpyxl.Workbook()
    wb.save(path)
    return path


# ---------------------------------------------------------------------------
# detect_file_type() tests
# ---------------------------------------------------------------------------

class TestDetectFileType:
    """Tests for file type detection with new formats."""

    def test_detect_pptx(self):
        assert detect_file_type("presentation.pptx") == "pptx"

    def test_detect_xlsx(self):
        assert detect_file_type("spreadsheet.xlsx") == "xlsx"

    def test_detect_pptx_case_insensitive(self):
        assert detect_file_type("SLIDES.PPTX") == "pptx"

    def test_detect_xlsx_case_insensitive(self):
        assert detect_file_type("DATA.XLSX") == "xlsx"

    def test_detect_unknown_extension(self):
        assert detect_file_type("archive.zip") == "unknown"

    def test_existing_types_still_work(self):
        """Ensure existing formats are not broken."""
        assert detect_file_type("doc.docx") == "docx"
        assert detect_file_type("file.pdf") == "pdf"
        assert detect_file_type("readme.txt") == "txt"
        assert detect_file_type("notes.md") == "md"


# ---------------------------------------------------------------------------
# extract_text_from_pptx() tests
# ---------------------------------------------------------------------------

class TestExtractTextFromPptx:
    """Tests for PowerPoint text extraction."""

    def test_single_slide_single_textbox(self, tmp_path):
        """Extract text from a single slide with one text box."""
        path = _create_pptx(str(tmp_path), [["Hello World"]])
        text = extract_text_from_pptx(path)
        assert "Hello World" in text
        assert "[Slide 1]" in text

    def test_single_slide_multiple_textboxes(self, tmp_path):
        """Extract text from multiple text boxes on one slide."""
        path = _create_pptx(str(tmp_path), [["Title", "Subtitle", "Body text"]])
        text = extract_text_from_pptx(path)
        assert "Title" in text
        assert "Subtitle" in text
        assert "Body text" in text

    def test_multiple_slides(self, tmp_path):
        """Text from multiple slides should be labeled with slide numbers."""
        path = _create_pptx(str(tmp_path), [
            ["Slide 1 content"],
            ["Slide 2 content"],
            ["Slide 3 content"],
        ])
        text = extract_text_from_pptx(path)
        assert "[Slide 1]" in text
        assert "[Slide 2]" in text
        assert "[Slide 3]" in text
        assert "Slide 1 content" in text
        assert "Slide 3 content" in text

    def test_table_in_slide(self, tmp_path):
        """Text from tables in slides should be extracted."""
        path = _create_pptx_with_table(str(tmp_path), [
            ["Header A", "Header B"],
            ["Value 1", "Value 2"],
        ])
        text = extract_text_from_pptx(path)
        assert "Header A" in text
        assert "Value 2" in text

    def test_empty_presentation(self, tmp_path):
        """Empty presentation should return empty string."""
        path = os.path.join(str(tmp_path), "empty.pptx")
        prs = Presentation()
        prs.save(path)
        text = extract_text_from_pptx(path)
        assert text == ""

    def test_slide_with_empty_textbox(self, tmp_path):
        """Slides with only empty text boxes should be skipped."""
        path = os.path.join(str(tmp_path), "empty_box.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
        txBox.text_frame.text = ""
        prs.save(path)
        text = extract_text_from_pptx(path)
        assert text == ""

    def test_chinese_content(self, tmp_path):
        """Chinese text should be handled correctly."""
        path = _create_pptx(str(tmp_path), [["研究方法", "實驗結果與討論"]])
        text = extract_text_from_pptx(path)
        assert "研究方法" in text
        assert "實驗結果與討論" in text


# ---------------------------------------------------------------------------
# extract_text_from_xlsx() tests
# ---------------------------------------------------------------------------

class TestExtractTextFromXlsx:
    """Tests for Excel text extraction."""

    def test_single_sheet_basic(self, tmp_path):
        """Extract text from a simple single-sheet workbook."""
        path = _create_xlsx(str(tmp_path), {
            "Data": [
                ["Name", "Age", "City"],
                ["Alice", 30, "Taipei"],
                ["Bob", 25, "Tokyo"],
            ],
        })
        text = extract_text_from_xlsx(path)
        assert "[Sheet: Data]" in text
        assert "Alice" in text
        assert "30" in text
        assert "Tokyo" in text

    def test_multiple_sheets(self, tmp_path):
        """Text from multiple sheets should include sheet names."""
        path = _create_xlsx(str(tmp_path), {
            "Students": [["Name", "GPA"], ["Alice", 3.9]],
            "Courses": [["Course", "Credits"], ["Math", 3]],
        })
        text = extract_text_from_xlsx(path)
        assert "[Sheet: Students]" in text
        assert "[Sheet: Courses]" in text
        assert "Alice" in text
        assert "Math" in text

    def test_empty_rows_skipped(self, tmp_path):
        """Completely empty rows should be skipped."""
        path = _create_xlsx(str(tmp_path), {
            "Sheet1": [
                ["Data"],
                [None, None, None],  # Empty row
                ["More data"],
            ],
        })
        text = extract_text_from_xlsx(path)
        lines = [l for l in text.split("\n") if l.strip() and not l.startswith("[")]
        assert len(lines) == 2  # "Data" and "More data"

    def test_numeric_values(self, tmp_path):
        """Numbers should be converted to strings."""
        path = _create_xlsx(str(tmp_path), {
            "Numbers": [[1, 2.5, 3.14159]],
        })
        text = extract_text_from_xlsx(path)
        assert "1" in text
        assert "2.5" in text
        assert "3.14159" in text

    def test_empty_workbook(self, tmp_path):
        """Workbook with only empty sheets should return empty string."""
        path = _create_xlsx_empty(str(tmp_path))
        text = extract_text_from_xlsx(path)
        assert text == ""

    def test_mixed_types(self, tmp_path):
        """Cells with different types should all be extracted."""
        path = _create_xlsx(str(tmp_path), {
            "Mixed": [
                ["Text", 42, True, None],
                ["More", 3.14, False, "End"],
            ],
        })
        text = extract_text_from_xlsx(path)
        assert "Text" in text
        assert "42" in text
        assert "True" in text
        assert "3.14" in text

    def test_chinese_content(self, tmp_path):
        """Chinese text should be handled correctly."""
        path = _create_xlsx(str(tmp_path), {
            "學生名單": [
                ["姓名", "學號", "系所"],
                ["王小明", "R12345", "資訊工程"],
            ],
        })
        text = extract_text_from_xlsx(path)
        assert "學生名單" in text
        assert "王小明" in text
        assert "資訊工程" in text

    def test_tab_separated_cells(self, tmp_path):
        """Cells should be tab-separated in output."""
        path = _create_xlsx(str(tmp_path), {
            "Data": [["A", "B", "C"]],
        })
        text = extract_text_from_xlsx(path)
        assert "A\tB\tC" in text


# ---------------------------------------------------------------------------
# extract_text() dispatcher tests (new formats)
# ---------------------------------------------------------------------------

class TestExtractTextDispatcher:
    """Tests for the unified extract_text() dispatcher with new formats."""

    def test_pptx_dispatch(self, tmp_path):
        """Dispatcher should route .pptx to extract_text_from_pptx."""
        path = _create_pptx(str(tmp_path), [["Dispatch test"]])
        text = extract_text(path, "pptx")
        assert "Dispatch test" in text

    def test_xlsx_dispatch(self, tmp_path):
        """Dispatcher should route .xlsx to extract_text_from_xlsx."""
        path = _create_xlsx(str(tmp_path), {"Sheet1": [["Dispatch test"]]})
        text = extract_text(path, "xlsx")
        assert "Dispatch test" in text

    def test_unsupported_type_raises(self, tmp_path):
        """Unsupported file type should raise ValueError."""
        with pytest.raises(ValueError, match="Unsupported file type"):
            extract_text("dummy.zip", "zip")

    def test_existing_docx_dispatch(self, tmp_path):
        """Existing docx dispatch should still work."""
        from docx import Document

        path = os.path.join(str(tmp_path), "test.docx")
        doc = Document()
        doc.add_paragraph("DOCX dispatch test")
        doc.save(path)
        text = extract_text(path, "docx")
        assert "DOCX dispatch test" in text

    def test_existing_txt_dispatch(self, tmp_path):
        """Existing txt dispatch should still work."""
        path = os.path.join(str(tmp_path), "test.txt")
        with open(path, "w", encoding="utf-8") as f:
            f.write("TXT dispatch test")
        text = extract_text(path, "txt")
        assert "TXT dispatch test" in text


# ---------------------------------------------------------------------------
# Config + supported extensions tests
# ---------------------------------------------------------------------------

class TestSupportedExtensions:
    """Tests for config supported extensions update."""

    def test_pptx_in_supported(self):
        """PPTX should be in supported extensions."""
        from app.config import settings
        exts = settings.get_supported_extensions()
        assert ".pptx" in exts

    def test_xlsx_in_supported(self):
        """XLSX should be in supported extensions."""
        from app.config import settings
        exts = settings.get_supported_extensions()
        assert ".xlsx" in exts

    def test_existing_extensions_preserved(self):
        """Existing extensions should still be present."""
        from app.config import settings
        exts = settings.get_supported_extensions()
        assert ".docx" in exts
        assert ".pdf" in exts
        assert ".txt" in exts
        assert ".md" in exts
