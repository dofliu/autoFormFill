import asyncio
import uuid

import fitz  # PyMuPDF
import openpyxl
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

from app.llm.factory import get_llm_adapter
from app.schemas.document import DocumentMetadataInput, DocumentUploadResponse
from app.utils.chunker import chunk_text
from app.vector_store import get_collection


def extract_text_from_docx(file_path: str) -> str:
    """Extract all text from a .docx file."""
    doc = DocxDocument(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text.strip())
    return "\n".join(paragraphs)


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF using pdfplumber (better for tables) with PyMuPDF fallback."""
    texts = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
    except Exception:
        # Fallback to PyMuPDF
        doc = fitz.open(file_path)
        for page in doc:
            texts.append(page.get_text())
        doc.close()
    return "\n".join(texts)


def extract_text_from_plaintext(file_path: str) -> str:
    """Read plain text files (.txt, .md) with encoding detection."""
    for encoding in ["utf-8", "utf-8-sig", "cp950", "big5", "latin-1"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    # Last resort: read as bytes and decode with replacement
    with open(file_path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint (.pptx) file.

    Reads text from all shapes (text boxes, titles, tables, grouped shapes)
    across every slide, preserving slide order.
    """
    prs = Presentation(file_path)
    texts: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            slide_texts.append(text)
        if slide_texts:
            texts.append(f"[Slide {slide_idx}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from an Excel (.xlsx) file.

    Reads cell values from all worksheets. Each sheet is prefixed with its name.
    Empty rows are skipped; cells are tab-separated.
    """
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    texts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_lines: list[str] = []
        for row in ws.iter_rows(values_only=True):
            # Filter out completely empty rows
            cell_values = [str(c) if c is not None else "" for c in row]
            if any(v.strip() for v in cell_values):
                sheet_lines.append("\t".join(cell_values))
        if sheet_lines:
            texts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_lines))
    wb.close()
    return "\n\n".join(texts)


def extract_text(file_path: str, file_type: str) -> str:
    """Extract text based on file type."""
    if file_type == "docx":
        return extract_text_from_docx(file_path)
    elif file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type in ("txt", "md"):
        return extract_text_from_plaintext(file_path)
    elif file_type == "pptx":
        return extract_text_from_pptx(file_path)
    elif file_type == "xlsx":
        return extract_text_from_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def _build_metadata(
    meta: DocumentMetadataInput,
    user_id: int | None = None,
) -> dict:
    """Build ChromaDB-compatible metadata dict (string values only).

    Args:
        meta: Document metadata from the upload request.
        user_id: Owner of the document. ``None`` → dev mode (no isolation).
    """
    result = {"title": meta.title, "doc_type": meta.doc_type}
    if meta.authors:
        result["authors"] = meta.authors
    if meta.publish_year:
        result["publish_year"] = str(meta.publish_year)
    if meta.keywords:
        result["keywords"] = meta.keywords
    if meta.project_name:
        result["project_name"] = meta.project_name
    if meta.funding_agency:
        result["funding_agency"] = meta.funding_agency
    if meta.execution_period:
        result["execution_period"] = meta.execution_period
    if meta.tech_stack:
        result["tech_stack"] = meta.tech_stack

    # Multi-user isolation metadata
    if user_id is not None:
        result["user_id"] = str(user_id)
        result["shared"] = "false"
    return result


async def embed_and_store(
    file_path: str,
    file_type: str,
    metadata: DocumentMetadataInput,
    user_id: int | None = None,
) -> DocumentUploadResponse:
    """Extract text, chunk, embed, and store in ChromaDB."""
    text = extract_text(file_path, file_type)
    if not text.strip():
        raise ValueError("No text could be extracted from the document.")

    chunks = chunk_text(text)
    doc_id = uuid.uuid4().hex

    # Get embeddings via LLM adapter
    adapter = get_llm_adapter()
    embeddings = await asyncio.to_thread(adapter.embed_batch, chunks)

    # Determine collection
    collection_name = (
        "academic_papers" if metadata.doc_type == "paper" else "research_projects"
    )
    collection = get_collection(collection_name)

    # Store in ChromaDB
    ids = [f"{doc_id}_{i}" for i in range(len(chunks))]
    meta_dict = _build_metadata(metadata, user_id=user_id)
    metadatas = [meta_dict for _ in chunks]

    await asyncio.to_thread(
        collection.add,
        ids=ids,
        documents=chunks,
        embeddings=embeddings,
        metadatas=metadatas,
    )

    return DocumentUploadResponse(
        doc_id=doc_id,
        collection=collection_name,
        chunks_count=len(chunks),
        metadata=meta_dict,
    )


async def search_documents(
    query: str,
    collection_name: str,
    n_results: int = 5,
    user_id: int | None = None,
) -> list[dict]:
    """Search ChromaDB for relevant document chunks.

    Args:
        query: Natural language search query.
        collection_name: Which ChromaDB collection to search.
        n_results: Maximum results to return.
        user_id: When set, only return documents owned by this user
                 **or** marked as shared. ``None`` → no filtering (dev mode).
    """
    adapter = get_llm_adapter()
    query_embedding = await asyncio.to_thread(adapter.embed_text, query)

    collection = get_collection(collection_name)

    # Build optional where filter for multi-user isolation
    where_filter = None
    if user_id is not None:
        where_filter = {
            "$or": [
                {"user_id": str(user_id)},
                {"shared": "true"},
            ]
        }

    results = await asyncio.to_thread(
        collection.query,
        query_embeddings=[query_embedding],
        n_results=n_results,
        include=["documents", "metadatas", "distances"],
        where=where_filter,
    )

    items = []
    if results and results["ids"] and results["ids"][0]:
        for i, doc_id in enumerate(results["ids"][0]):
            items.append(
                {
                    "doc_id": doc_id,
                    "text": results["documents"][0][i],
                    "metadata": results["metadatas"][0][i],
                    "distance": results["distances"][0][i] if results["distances"] else None,
                }
            )
    return items
